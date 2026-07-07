from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path(__file__).resolve().parent
MD_PATH = BASE / "while_python_youth_lesson.md"
PPTX_PATH = BASE / "while_python_youth_lesson.pptx"

BG = RGBColor(12, 14, 24)
ACCENT = RGBColor(0, 255, 170)
PINK = RGBColor(255, 64, 129)
TEXT = RGBColor(245, 247, 250)
MUTED = RGBColor(180, 190, 205)
CODE_BG = RGBColor(28, 32, 48)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_textbox(slide, x, y, w, h, font_size=24, color=TEXT, bold=False):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.font.name = "Arial"
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    return shape


def add_tag(slide, text, x, y, w, color=ACCENT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.34))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = RGBColor(10, 10, 14)


def parse_slide(part):
    lines = [line.rstrip() for line in part.splitlines()]
    heading = ""
    body = []
    code_blocks = []
    in_code = False
    current_code = []

    for line in lines:
        if line.startswith("```"):
            if in_code:
                code_blocks.append("\n".join(current_code).strip())
                current_code = []
                in_code = False
            else:
                in_code = True
            continue

        if in_code:
            current_code.append(line)
            continue

        if line.startswith("# "):
            heading = line[2:].strip()
        elif line.startswith("## ") and not heading:
            heading = line[3:].strip()
        elif line.strip():
            cleaned = line.strip()
            if cleaned.startswith("## "):
                cleaned = cleaned[3:].strip()
            body.append(cleaned)

    if not heading:
        heading = body.pop(0) if body else "while"

    return heading, body, code_blocks


def add_title(slide, title):
    shape = add_textbox(slide, Inches(0.65), Inches(0.42), Inches(12), Inches(0.9), 32, TEXT, True)
    shape.text_frame.paragraphs[0].text = title[:95]


def add_body(slide, body, code_blocks):
    width = Inches(7.15 if code_blocks else 11.9)
    box = add_textbox(slide, Inches(0.75), Inches(1.55), width, Inches(5.55), 20, TEXT)
    tf = box.text_frame
    tf.clear()

    for index, item in enumerate(body[:12]):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        text = item.replace("`", "")

        if text.startswith("- "):
            p.text = text[2:]
            p.font.size = Pt(19)
            p.font.color.rgb = TEXT
        elif text.lower().startswith("мем на слайд"):
            p.text = "MEM: " + text.split(":", 1)[-1].strip()
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = PINK
        elif text.startswith("> "):
            p.text = text[2:]
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = ACCENT
        else:
            p.text = text
            p.font.size = Pt(18)
            p.font.color.rgb = MUTED if text.endswith(":") else TEXT
            p.font.bold = text.endswith(":")

        p.font.name = "Arial"

    if code_blocks:
        code = code_blocks[0]
        if len(code) > 650:
            code = code[:650] + "\n# ..."

        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.1),
            Inches(1.65),
            Inches(4.55),
            Inches(4.95),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = CODE_BG
        panel.line.color.rgb = ACCENT
        panel.line.width = Pt(1.5)

        tf = panel.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Pt(14)
        tf.margin_right = Pt(10)
        tf.margin_top = Pt(12)
        p = tf.paragraphs[0]
        p.text = code
        p.font.name = "Consolas"
        p.font.size = Pt(12 if len(code.splitlines()) > 13 else 14)
        p.font.color.rgb = RGBColor(225, 235, 255)


def add_footer(slide, index):
    add_tag(slide, f"WHILE #{index}", Inches(0.72), Inches(6.95), Inches(1.25))
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.1), Inches(7.08), Inches(10.5), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(55, 65, 90)
    line.line.fill.background()


def main():
    parts = [part.strip() for part in MD_PATH.read_text(encoding="utf-8").split("---") if part.strip()]
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for index, part in enumerate(parts, start=1):
        title, body, code_blocks = parse_slide(part)
        slide = prs.slides.add_slide(blank)
        set_bg(slide)
        add_title(slide, title)
        add_body(slide, body, code_blocks)
        add_footer(slide, index)

    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
